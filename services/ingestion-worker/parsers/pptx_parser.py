"""PPTX parser using python-pptx. Returns one text block per slide."""
from pptx import Presentation
from pptx.util import Pt


def parse_pptx(path: str) -> tuple[list[str], bool]:
    """Returns (slides: list[str], ocr_used=False)."""
    prs = Presentation(path)
    slides: list[str] = []

    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs if run.text).strip()
                if line:
                    parts.append(line)
        slides.append("\n".join(parts))

    return slides, False
