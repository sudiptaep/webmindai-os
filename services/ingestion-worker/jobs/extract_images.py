"""
Image extraction — F-17-A.

Extracts images from PDFs (PyMuPDF) and PPTX (python-pptx) via two strategies:
  1. Embedded XObject bitmaps — via page.get_images() (photos, scanned images)
  2. Page-render fallback — renders visually rich pages (flowcharts, anatomical
     diagrams, circuit diagrams drawn as vectors) as pixmaps at 150 DPI.
     Triggered when:
       a. Page has >= PAGE_RENDER_MIN_DRAWINGS vector paths (general heuristic), OR
       b. Page has a figure caption ("Figure X:", "Fig. X:") AND >= 3 vector paths
          (figure-caption-guided targeting — catches sparse diagrams missed by (a))
     This captures the majority of educational diagram content that get_images() misses.

Figure caption mining: scan each page's text for "Figure X:" / "Fig. X:" patterns.
Ground-truth captions are attached to extracted image records so the vision model
can use the exact caption instead of guessing — significantly better RAG indexing.

Filters out icons/decorations, dedups by content hash, writes JPEG + thumbnail.
"""
import hashlib
import io
import logging
import os
import re

import fitz  # PyMuPDF
from PIL import Image

logger = logging.getLogger(__name__)

IMAGE_MIN_WIDTH_PX = int(os.environ.get("IMAGE_MIN_WIDTH_PX", "100"))
IMAGE_MIN_HEIGHT_PX = int(os.environ.get("IMAGE_MIN_HEIGHT_PX", "100"))
IMAGE_MIN_SIZE_BYTES = int(os.environ.get("IMAGE_MIN_SIZE_BYTES", "10000"))
IMAGE_MAX_ASPECT_RATIO = float(os.environ.get("IMAGE_MAX_ASPECT_RATIO", "15"))
IMAGE_JPEG_QUALITY = int(os.environ.get("IMAGE_JPEG_QUALITY", "85"))
IMAGE_THUMBNAIL_SIZE = int(os.environ.get("IMAGE_THUMBNAIL_SIZE", "200"))
IMAGE_THUMBNAIL_QUALITY = int(os.environ.get("IMAGE_THUMBNAIL_QUALITY", "75"))

# Page render settings
PAGE_RENDER_DPI = int(os.environ.get("PAGE_RENDER_DPI", "150"))
# Min drawing paths for a page to be considered visually rich (general heuristic)
PAGE_RENDER_MIN_DRAWINGS = int(os.environ.get("PAGE_RENDER_MIN_DRAWINGS", "20"))
# Min drawings when a page also has an explicit figure caption (much lower threshold —
# sparse diagrams like single anatomical cross-sections may have only 3–12 paths)
PAGE_RENDER_MIN_DRAWINGS_WITH_CAPTION = int(os.environ.get("PAGE_RENDER_MIN_DRAWINGS_WITH_CAPTION", "8"))
# Max fraction of page area covered by text blocks to still render as image
# (a page that is >90% text is just a text page, not a diagram)
PAGE_RENDER_MAX_TEXT_RATIO = float(os.environ.get("PAGE_RENDER_MAX_TEXT_RATIO", "0.80"))
# Figure-region cropping: the vector drawing cluster's bounding box must cover at least
# this fraction of the page to count as a real figure (excludes stray marks/underlines).
PAGE_RENDER_MIN_REGION_RATIO = float(os.environ.get("PAGE_RENDER_MIN_REGION_RATIO", "0.12"))
# If the drawing cluster covers more than this fraction, treat as a full-page diagram.
PAGE_RENDER_FULLPAGE_REGION_RATIO = float(os.environ.get("PAGE_RENDER_FULLPAGE_REGION_RATIO", "0.85"))
# A single drawing rect larger than this fraction of the page is a border/background — ignore it.
PAGE_RENDER_BORDER_RECT_RATIO = float(os.environ.get("PAGE_RENDER_BORDER_RECT_RATIO", "0.9"))
# Padding (PDF points) added around the detected figure region before rendering.
PAGE_RENDER_REGION_PADDING_PT = float(os.environ.get("PAGE_RENDER_REGION_PADDING_PT", "10"))
# Max vertical gap (PDF points) between a figure-caption text block and the drawing
# cluster for the caption to count as "belonging" to that figure. This distinguishes
# a page that CONTAINS a figure (caption sits right below/above the drawing) from one
# that merely MENTIONS a figure in body prose (caption text far from any drawing cluster).
CAPTION_ADJACENCY_PT = float(os.environ.get("CAPTION_ADJACENCY_PT", "35"))
# Proximity gap (PDF points) for grouping drawing rects into a cluster. Elements of one
# figure sit within this gap; scattered page furniture (header rules, margin marks) does
# not merge, so the largest cluster is the actual figure — not a full-page union.
CLUSTER_GAP_PT = float(os.environ.get("CLUSTER_GAP_PT", "18"))

# Figure caption pattern: matches "Figure 1-2:", "Fig. 3A:", "Figure 6-11. Description"
# Captures: (label, figure_number, caption_text)
# Separators: colon, em-dash, or period (Guyton and many medical texts use period).
# Word-boundary \b prevents matching inside longer words; figure number must be
# followed immediately by separator — this filters most mid-sentence references
# like "as shown in Figure 3-1, the..." (which have comma/space/paren, not separator).
_FIG_CAPTION_RE = re.compile(
    r'\b(Fig(?:ure)?|FIGURE|FIG\.?|Plate|PLATE)\s*\.?\s*(\d+[\.\-]?\d*[A-Za-z]?)'
    r'\s*[:\.\—–]\s*([^\n]{10,250})',
    re.IGNORECASE,
)


def _extract_figure_captions(pdf: "fitz.Document") -> dict[int, list[dict]]:
    """
    Scan all pages for figure caption text ("Figure X: description").
    Returns {1-based page_num: [{fig_label, fig_num, caption_text}, ...]}.
    Also checks adjacent page (+1) so captions that wrap to the next page
    get associated with the figure page.
    """
    raw: dict[int, list[dict]] = {}
    for page_num in range(len(pdf)):
        text = pdf[page_num].get_text("text")
        matches = _FIG_CAPTION_RE.findall(text)
        if matches:
            raw[page_num + 1] = [
                {
                    "fig_label": m[0],
                    "fig_num": m[1],
                    "caption_text": m[2].strip()[:250],
                }
                for m in matches
            ]
    # Propagate: if page N+1 has a caption, also attach it to page N
    # (caption below a figure can be the first text on the next page in some PDFs)
    result: dict[int, list[dict]] = {}
    for page_num, caps in raw.items():
        result.setdefault(page_num, []).extend(caps)
        prev = page_num - 1
        if prev >= 1:
            result.setdefault(prev, []).extend(caps)
    return result


def _images_dir(college_id: str, doc_id: str) -> str:
    storage_root = (
        os.environ.get("STORAGE_ROOT")
        or os.environ.get("UPLOADS_DIR")
        or os.path.join(os.getcwd(), "uploads")
    )
    path = os.path.join(storage_root, "colleges", college_id, "images", doc_id)
    os.makedirs(path, exist_ok=True)
    return path


def _content_hash(img_bytes: bytes) -> str:
    return hashlib.md5(img_bytes).hexdigest()


def _should_skip(width: int, height: int, size_bytes: int) -> str | None:
    """Returns a filter_reason string if the image should be rejected, else None."""
    if width < IMAGE_MIN_WIDTH_PX or height < IMAGE_MIN_HEIGHT_PX:
        return "too_small"
    if size_bytes < IMAGE_MIN_SIZE_BYTES:
        return "too_small"
    if width > 0 and height > 0:
        ratio = max(width, height) / min(width, height)
        if ratio > IMAGE_MAX_ASPECT_RATIO:
            return "logo_icon"
    return None


def _largest_cluster_bbox(rects: "list[fitz.Rect]", gap: float) -> "fitz.Rect":
    """
    Group rects into clusters where members sit within `gap` points of each other,
    then return the bounding box of the cluster with the largest area. Iterates until
    the merge set is stable so proximity is transitive (A near B near C → one cluster).
    """
    boxes = [fitz.Rect(r) for r in rects]
    merged = True
    while merged and len(boxes) > 1:
        merged = False
        out: list[fitz.Rect] = []
        for r in boxes:
            grown = fitz.Rect(r.x0 - gap, r.y0 - gap, r.x1 + gap, r.y1 + gap)
            placed = False
            for c in out:
                if c.intersects(grown):
                    c.include_rect(r)
                    placed = True
                    merged = True
                    break
            if not placed:
                out.append(fitz.Rect(r))
        boxes = out
    return max(boxes, key=lambda b: b.width * b.height)


def _caption_adjacent_to_region(page: "fitz.Page", region: "fitz.Rect") -> bool:
    """
    True if a figure-caption text block ("Figure X: …") sits within
    CAPTION_ADJACENCY_PT (vertically) of the drawing cluster and overlaps it
    horizontally — i.e. the caption belongs to this figure. A prose mention of a
    figure elsewhere on the page is far from the cluster, so it returns False.
    """
    for b in page.get_text("blocks"):
        if len(b) < 7 or b[6] != 0:  # type 0 = text block
            continue
        block_text = b[4] if len(b) > 4 else ""
        if not _FIG_CAPTION_RE.search(block_text):
            continue
        bx0, by0, bx1, by1 = b[0], b[1], b[2], b[3]
        vgap = max(region.y0 - by1, by0 - region.y1, 0.0)  # 0 if overlapping vertically
        h_overlap = min(region.x1, bx1) - max(region.x0, bx0)
        if vgap <= CAPTION_ADJACENCY_PT and h_overlap > 0:
            return True
    return False


def _figure_region_bbox(page: "fitz.Page", has_fig_caption: bool = False) -> "fitz.Rect | None":
    """
    Compute the bounding box of the vector-drawing cluster on a page — the actual
    figure region — so we render ONLY the figure, not the whole page.

    A page is treated as containing a figure when EITHER:
      - it has strong vector content (>= PAGE_RENDER_MIN_DRAWINGS paths), OR
      - a figure caption text block sits adjacent to the drawing cluster.
    Pages that merely mention "Figure X" in body prose (caption far from any cluster,
    little vector content) return None and are skipped — this is what prevents the
    "render every page of the textbook" blow-up.

    Returns a padded fitz.Rect clipped to the page, or None.
    """
    drawings = page.get_drawings()
    # Need at least a few drawing paths to bother; the real gate is below.
    if len(drawings) < PAGE_RENDER_MIN_DRAWINGS_WITH_CAPTION:
        return None

    page_rect = page.rect
    page_area = page_rect.width * page_rect.height
    if page_area <= 0:
        return None

    rects: list[fitz.Rect] = []
    for d in drawings:
        r = d.get("rect")
        if r is None:
            continue
        if r.width <= 1 or r.height <= 1:  # hairline rule / degenerate stroke
            continue
        if (r.width * r.height) >= PAGE_RENDER_BORDER_RECT_RATIO * page_area:
            continue  # page border / full-page background box
        rects.append(r)

    if not rects:
        return None

    # Group rects into spatial clusters and take the largest — isolates the dense
    # figure from scattered page furniture (which would otherwise inflate a global
    # union bbox to the whole page).
    union = _largest_cluster_bbox(rects, CLUSTER_GAP_PT)
    union.intersect(page_rect)
    if union.is_empty or union.width <= 0 or union.height <= 0:
        return None

    ratio = (union.width * union.height) / page_area
    if ratio >= PAGE_RENDER_FULLPAGE_REGION_RATIO:
        union = fitz.Rect(page_rect)  # whole-page diagram
    elif ratio < PAGE_RENDER_MIN_REGION_RATIO:
        return None  # cluster too small — stray marks, not a figure

    # Real-figure decision: strong vector content OR a caption adjacent to the cluster.
    strong_vector = len(drawings) >= PAGE_RENDER_MIN_DRAWINGS
    if not strong_vector:
        if not (has_fig_caption and _caption_adjacent_to_region(page, union)):
            return None

    pad = PAGE_RENDER_REGION_PADDING_PT
    padded = fitz.Rect(union.x0 - pad, union.y0 - pad, union.x1 + pad, union.y1 + pad)
    padded.intersect(page_rect)
    return padded


def _render_page_region(
    page: "fitz.Page", clip: "fitz.Rect", page_num: int, images_dir: str, img_idx: int
) -> dict | None:
    """Renders a clipped figure region of a PDF page at PAGE_RENDER_DPI as JPEG + thumbnail."""
    try:
        mat = fitz.Matrix(PAGE_RENDER_DPI / 72, PAGE_RENDER_DPI / 72)
        pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
        img_bytes = pix.tobytes("jpeg")

        pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        filter_reason = _should_skip(pil_img.width, pil_img.height, len(img_bytes))
        if filter_reason:
            return None

        filename = f"page{page_num:04d}_fig{img_idx:02d}.jpg"
        file_path_out = os.path.join(images_dir, filename)
        pil_img.save(file_path_out, "JPEG", quality=IMAGE_JPEG_QUALITY)

        thumb_path = os.path.join(images_dir, f"thumb_{filename}")
        thumb = pil_img.copy()
        thumb.thumbnail((IMAGE_THUMBNAIL_SIZE, IMAGE_THUMBNAIL_SIZE), Image.LANCZOS)
        thumb.save(thumb_path, "JPEG", quality=IMAGE_THUMBNAIL_QUALITY)

        return {
            "was_filtered": False,
            "file_path": file_path_out,
            "thumbnail_path": thumb_path,
            "file_size_bytes": os.path.getsize(file_path_out),
            "width_px": pil_img.width,
            "height_px": pil_img.height,
            "format": "jpg",
            "content_hash": _content_hash(img_bytes),
            "source_page": page_num,
            "image_index_on_page": img_idx,
            "global_image_index": -1,  # caller assigns final index
        }
    except Exception as exc:
        logger.warning("Page region render failed page=%d: %s", page_num, exc)
        return None


def extract_images_from_pdf(file_path: str, doc_id: str, college_id: str) -> list[dict]:
    """
    Extract images from a PDF using two strategies:
    1. Embedded XObject bitmaps via get_images() — photos, scanned images.
    2. Page-render fallback — renders visually rich pages (vector diagrams, flowcharts,
       anatomical figures) that get_images() completely misses.
       Triggered by either:
         a. >= PAGE_RENDER_MIN_DRAWINGS vector paths (general heuristic)
         b. Page has an explicit "Figure X:" caption AND >= 3 paths (caption-guided)

    Figure captions found in page text are attached to each image record as
    `figure_captions` so the vision model can use the ground-truth caption text.

    Returns a list of image metadata dicts (was_filtered=True entries included).
    """
    images_dir = _images_dir(college_id, doc_id)
    pdf = fitz.open(file_path)

    # Pre-scan all figure captions once — O(n_pages) text scan, fast
    figure_captions = _extract_figure_captions(pdf)
    logger.info("PDF figure caption scan: %d pages have Fig captions", len(figure_captions))

    extracted: list[dict] = []
    seen_hashes: set[str] = set()
    global_index = 0

    for page_num in range(len(pdf)):
        page = pdf[page_num]
        image_list = page.get_images(full=True)
        page_qualified_count = 0  # how many non-filtered bitmaps this page produced
        page_captions = figure_captions.get(page_num + 1, [])

        # Strategy 1: embedded bitmaps
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_image = pdf.extract_image(xref)
                img_bytes = base_image["image"]
                img_width = base_image["width"]
                img_height = base_image["height"]

                filter_reason = _should_skip(img_width, img_height, len(img_bytes))
                if filter_reason:
                    extracted.append({
                        "was_filtered": True,
                        "filter_reason": filter_reason,
                        "source_page": page_num + 1,
                    })
                    continue

                content_hash = _content_hash(img_bytes)
                if content_hash in seen_hashes:
                    extracted.append({
                        "was_filtered": True,
                        "filter_reason": "duplicate",
                        "source_page": page_num + 1,
                        "content_hash": content_hash,
                    })
                    continue
                seen_hashes.add(content_hash)

                filename = f"page{page_num + 1:04d}_img{img_idx:02d}.jpg"
                file_path_out = os.path.join(images_dir, filename)

                pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                pil_img.save(file_path_out, "JPEG", quality=IMAGE_JPEG_QUALITY)

                thumb_path = os.path.join(images_dir, f"thumb_{filename}")
                thumb = pil_img.copy()
                thumb.thumbnail((IMAGE_THUMBNAIL_SIZE, IMAGE_THUMBNAIL_SIZE), Image.LANCZOS)
                thumb.save(thumb_path, "JPEG", quality=IMAGE_THUMBNAIL_QUALITY)

                extracted.append({
                    "was_filtered": False,
                    "file_path": file_path_out,
                    "thumbnail_path": thumb_path,
                    "file_size_bytes": os.path.getsize(file_path_out),
                    "width_px": pil_img.width,
                    "height_px": pil_img.height,
                    "format": "jpg",
                    "content_hash": content_hash,
                    "source_page": page_num + 1,
                    "image_index_on_page": img_idx,
                    "global_image_index": global_index,
                    "figure_captions": page_captions,
                })
                global_index += 1
                page_qualified_count += 1

            except Exception as exc:
                logger.warning("Skip image xref=%s on page %d: %s", xref, page_num + 1, exc)
                continue

        # Strategy 2: figure-region render for vector-drawn diagrams with no qualifying bitmaps.
        # Catches flowcharts, anatomical diagrams, circuit diagrams, LaTeX figures — all
        # stored as PDF vector paths, invisible to get_images(). We render ONLY the drawing
        # cluster's bounding box (not the whole page), so body text is excluded and pure-text
        # pages that merely mention "Figure X" in prose produce no cluster → are skipped.
        # Fig-caption-guided: an explicit caption lowers the drawing-count threshold.
        if page_qualified_count == 0:
            region = _figure_region_bbox(page, has_fig_caption=bool(page_captions))
            if region is not None:
                rendered = _render_page_region(page, region, page_num + 1, images_dir, img_idx=0)
                if rendered is not None:
                    content_hash = rendered["content_hash"]
                    if content_hash not in seen_hashes:
                        seen_hashes.add(content_hash)
                        rendered["global_image_index"] = global_index
                        rendered["figure_captions"] = page_captions
                        extracted.append(rendered)
                        global_index += 1
                        logger.debug(
                            "Page %d figure region rendered (captions=%d)",
                            page_num + 1, len(page_captions),
                        )

    pdf.close()
    return extracted


def extract_images_from_pptx(file_path: str, doc_id: str, college_id: str) -> list[dict]:
    """
    Extract all embedded images from a PPTX file. Each slide is treated as a "page".
    """
    from pptx import Presentation

    MSO_PICTURE = 13

    images_dir = _images_dir(college_id, doc_id)
    prs = Presentation(file_path)
    extracted: list[dict] = []
    seen_hashes: set[str] = set()
    global_index = 0

    for slide_num, slide in enumerate(prs.slides):
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_PICTURE:
                continue
            try:
                img_blob = shape.image.blob
                pil_img = Image.open(io.BytesIO(img_blob)).convert("RGB")

                filter_reason = _should_skip(pil_img.width, pil_img.height, len(img_blob))
                if filter_reason:
                    extracted.append({
                        "was_filtered": True,
                        "filter_reason": filter_reason,
                        "source_page": slide_num + 1,
                    })
                    img_idx += 1
                    continue

                content_hash = _content_hash(img_blob)
                if content_hash in seen_hashes:
                    extracted.append({
                        "was_filtered": True,
                        "filter_reason": "duplicate",
                        "source_page": slide_num + 1,
                        "content_hash": content_hash,
                    })
                    img_idx += 1
                    continue
                seen_hashes.add(content_hash)

                filename = f"slide{slide_num + 1:04d}_img{img_idx:02d}.jpg"
                file_path_out = os.path.join(images_dir, filename)
                pil_img.save(file_path_out, "JPEG", quality=IMAGE_JPEG_QUALITY)

                thumb_path = os.path.join(images_dir, f"thumb_{filename}")
                thumb = pil_img.copy()
                thumb.thumbnail((IMAGE_THUMBNAIL_SIZE, IMAGE_THUMBNAIL_SIZE), Image.LANCZOS)
                thumb.save(thumb_path, "JPEG", quality=IMAGE_THUMBNAIL_QUALITY)

                extracted.append({
                    "was_filtered": False,
                    "file_path": file_path_out,
                    "thumbnail_path": thumb_path,
                    "file_size_bytes": os.path.getsize(file_path_out),
                    "width_px": pil_img.width,
                    "height_px": pil_img.height,
                    "format": "jpg",
                    "content_hash": content_hash,
                    "source_page": slide_num + 1,
                    "image_index_on_page": img_idx,
                    "global_image_index": global_index,
                })
                img_idx += 1
                global_index += 1

            except Exception as exc:
                logger.warning("Skip PPTX image slide=%d: %s", slide_num + 1, exc)
                img_idx += 1
                continue

    return extracted
